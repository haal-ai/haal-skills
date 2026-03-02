#!/usr/bin/env python3
"""Generate a PowerPoint presentation from a structured markdown plan."""

import sys
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Color palette ---
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_SLIDE = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_PURPLE = RGBColor(0x7B, 0x2F, 0xFF)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_MUTED = RGBColor(0x99, 0x99, 0x99)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
TABLE_HEADER_BG = RGBColor(0x0F, 0x3D, 0x6B)
TABLE_ROW_BG = RGBColor(0x12, 0x2A, 0x4F)
TABLE_ALT_BG = RGBColor(0x0E, 0x1F, 0x3A)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def parse_markdown(filepath):
    """Parse the markdown plan into slides."""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # Extract title
    title_match = re.match(r'^#\s+(.+)', content)
    title = title_match.group(1).strip() if title_match else "Presentation"

    # Split by slide headers
    slide_pattern = r'##\s+Slide\s+\d+:\s+'
    parts = re.split(slide_pattern, content)
    headers = re.findall(r'##\s+Slide\s+\d+:\s+(.+)', content)

    slides = []
    for i, header in enumerate(headers):
        body = parts[i + 1].strip()
        # Remove trailing ---
        body = re.sub(r'\n---\s*$', '', body).strip()
        slides.append({"title": header.strip(), "body": body})

    return title, slides


def parse_body_elements(body):
    """Parse body into structured elements: subtitle, paragraphs, bullets, tables, code blocks."""
    elements = []
    lines = body.split('\n')
    i = 0

    # Check for subtitle (bold line at start)
    if lines and re.match(r'^\*\*(.+)\*\*$', lines[0].strip()):
        elements.append({"type": "subtitle", "text": re.match(r'^\*\*(.+)\*\*$', lines[0].strip()).group(1)})
        i = 1

    while i < len(lines):
        line = lines[i]

        # Code block
        if line.strip().startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            elements.append({"type": "code", "text": '\n'.join(code_lines)})
            i += 1
            continue

        # Table
        if '|' in line and i + 1 < len(lines) and re.match(r'^[\s|:-]+$', lines[i + 1]):
            headers_raw = [c.strip() for c in line.strip().strip('|').split('|')]
            headers_clean = [re.sub(r'\*\*(.+?)\*\*', r'\1', h) for h in headers_raw]
            i += 2  # skip header + separator
            rows = []
            while i < len(lines) and '|' in lines[i] and lines[i].strip():
                cells = [c.strip() for c in lines[i].strip().strip('|').split('|')]
                cells = [re.sub(r'\*\*(.+?)\*\*', r'\1', c) for c in cells]
                rows.append(cells)
                i += 1
            elements.append({"type": "table", "headers": headers_clean, "rows": rows})
            continue

        # Bullet point
        if re.match(r'^-\s+', line.strip()):
            text = re.sub(r'^-\s+', '', line.strip())
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            elements.append({"type": "bullet", "text": text})
            i += 1
            continue

        # Paragraph text
        stripped = line.strip()
        if stripped and not stripped.startswith('#'):
            stripped = re.sub(r'\*\*(.+?)\*\*', r'\1', stripped)
            elements.append({"type": "paragraph", "text": stripped})

        i += 1

    return elements



def build_title_slide(prs, title):
    """Create the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Accent bar top
    add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

    # Title
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5),
                 title, font_size=36, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtle bottom bar
    add_accent_bar(slide, Inches(4), Inches(4.0), Inches(5), Inches(0.04), ACCENT_PURPLE)


def render_table(slide, element, top_pos):
    """Render a table element on the slide and return the bottom position."""
    headers = element["headers"]
    rows = element["rows"]
    cols = len(headers)
    row_count = len(rows) + 1

    table_width = Inches(11.0)
    col_width = table_width // cols
    row_height = Inches(0.45)
    table_height = row_height * row_count

    left = Inches(1.0)
    table_shape = slide.shapes.add_table(row_count, cols, left, top_pos, table_width, table_height)
    table = table_shape.table

    # Style header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = ACCENT_BLUE
            paragraph.font.name = "Segoe UI"
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    # Style data rows
    for ri, row in enumerate(rows):
        for ci in range(cols):
            cell = table.cell(ri + 1, ci)
            cell.text = row[ci] if ci < len(row) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = TEXT_LIGHT
                paragraph.font.name = "Segoe UI"
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_BG if ri % 2 == 0 else TABLE_ALT_BG

    return top_pos + table_height + Inches(0.2)


def build_content_slide(prs, slide_data, slide_index):
    """Create a content slide from parsed data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_SLIDE)

    # Top accent bar
    add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.04), ACCENT_BLUE)

    # Slide number
    add_text_box(slide, Inches(11.5), Inches(0.15), Inches(1.2), Inches(0.3),
                 f"{slide_index}", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.6),
                 slide_data["title"], font_size=26, color=TEXT_WHITE, bold=True)

    # Accent bar under title
    add_accent_bar(slide, Inches(0.8), Inches(0.95), Inches(2.0), Inches(0.04), ACCENT_BLUE)

    elements = parse_body_elements(slide_data["body"])
    top = Inches(1.2)

    for elem in elements:
        if top > Inches(6.8):
            break  # prevent overflow

        if elem["type"] == "subtitle":
            add_text_box(slide, Inches(0.8), top, Inches(11.5), Inches(0.4),
                         elem["text"], font_size=16, color=ACCENT_BLUE, bold=False)
            top += Inches(0.45)

        elif elem["type"] == "paragraph":
            text = elem["text"]
            line_count = max(1, len(text) // 100 + 1)
            height = Inches(0.3 * line_count)
            add_text_box(slide, Inches(0.8), top, Inches(11.5), height,
                         text, font_size=13, color=TEXT_LIGHT)
            top += height + Inches(0.1)

        elif elem["type"] == "bullet":
            text = elem["text"]
            line_count = max(1, len(text) // 90 + 1)
            height = Inches(0.28 * line_count)
            add_text_box(slide, Inches(1.2), top, Inches(11.0), height,
                         f"▸  {text}", font_size=12, color=TEXT_LIGHT)
            top += height + Inches(0.06)

        elif elem["type"] == "table":
            top = render_table(slide, elem, top)

        elif elem["type"] == "code":
            code_text = elem["text"]
            line_count = code_text.count('\n') + 1
            height = Inches(0.22 * min(line_count, 8))
            # Code background box
            code_bg = slide.shapes.add_shape(1, Inches(0.8), top, Inches(11.5), height + Inches(0.2))
            code_bg.fill.solid()
            code_bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x1A)
            code_bg.line.fill.background()
            add_text_box(slide, Inches(1.0), top + Inches(0.1), Inches(11.0), height,
                         code_text, font_size=9, color=ACCENT_GREEN, font_name="Consolas")
            top += height + Inches(0.35)


def generate_pptx(plan_path, output_dir):
    """Main generation function."""
    title, slides = parse_markdown(plan_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    build_title_slide(prs, title)

    # Content slides
    for i, slide_data in enumerate(slides):
        build_content_slide(prs, slide_data, i)

    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)

    # Build filename
    plan_name = os.path.splitext(os.path.basename(plan_path))[0]
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d-%H%M")
    output_path = os.path.join(output_dir, f"{plan_name}-{timestamp}.pptx")

    prs.save(output_path)
    print(f"PowerPoint generated: {output_path}")
    print(f"Slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_dynamic_pptx.py <plan.md> [output_dir]")
        sys.exit(1)

    plan_file = sys.argv[1]
    out_dir = sys.argv[2] if len(sys.argv) > 2 else ".olaf/work/staging/pptx-folder/"
    generate_pptx(plan_file, out_dir)
