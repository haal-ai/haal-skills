#!/usr/bin/env python3
"""
Dynamic PowerPoint Generation from Plan
=======================================
Generates professional PowerPoint presentations from markdown presentation plans.
Uses modern layout resolution and theme-aware styling while maintaining flexibility
for dynamic content creation.

THEME-AWARE: All colors use PowerPoint theme slots so the user can switch
the theme in PowerPoint and all colors adapt automatically.

LAYOUT-AWARE: Slide layouts are resolved by name from the template, not by
hardcoded index. Titles are populated via the layout's built-in placeholders.

FLEXIBLE: Creates slides dynamically based on plan content without rigid
constraints on slide count or format.

Requirements: pip install python-pptx
"""

import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime


# ── Theme color semantic names ──
ACCENT_PRIMARY   = MSO_THEME_COLOR.ACCENT_1   # Blue / structural
ACCENT_POSITIVE  = MSO_THEME_COLOR.ACCENT_2   # Green / output / quality
ACCENT_AI        = MSO_THEME_COLOR.ACCENT_3   # Purple / validation
ACCENT_WARNING   = MSO_THEME_COLOR.ACCENT_4   # Orange / config
ACCENT_DANGER    = MSO_THEME_COLOR.ACCENT_5   # Red / problems
ACCENT_INFO      = MSO_THEME_COLOR.ACCENT_6   # Teal / processing


# ── Layout resolution ──
# Layouts are resolved by name (case-insensitive substring match).
# Customize these patterns if your template uses different naming.
LAYOUT_COVER_NAMES = ["cover", "title slide"]
LAYOUT_TITLE_ONLY_NAMES = ["title only", "title-only"]
LAYOUT_CLOSING_NAMES = ["closing", "end", "thank"]


def _find_layout(prs, name_patterns, fallback_index=0):
    """Find a slide layout by name pattern (case-insensitive substring match).
    Tries each pattern in order and returns the first matching layout.
    Falls back to the given index if no name matches.
    """
    for pattern in name_patterns:
        for layout in prs.slide_layouts:
            if pattern.lower() in layout.name.lower():
                return layout
    if fallback_index < len(prs.slide_layouts):
        return prs.slide_layouts[fallback_index]
    return prs.slide_layouts[0]


def get_cover_layout(prs):
    """Get the cover/title slide layout."""
    return _find_layout(prs, LAYOUT_COVER_NAMES, fallback_index=0)


def get_title_only_layout(prs):
    """Get the title-only slide layout (REQUIRED for all non-cover slides)."""
    return _find_layout(prs, LAYOUT_TITLE_ONLY_NAMES, fallback_index=1)


def get_closing_layout(prs):
    """Get the closing/outro slide layout."""
    return _find_layout(prs, LAYOUT_CLOSING_NAMES, fallback_index=0)


# ── Helpers ──

def _apply_theme_color(color_format, theme_color, brightness=0.0):
    """Apply a theme color with optional brightness adjustment."""
    color_format.theme_color = theme_color
    if brightness != 0.0:
        color_format.brightness = brightness


def set_bg(slide):
    """No-op: let the slide inherit its background from the slide master/layout."""
    pass


def txt(slide, l, t, w, h, text, sz=14, theme_color=MSO_THEME_COLOR.LIGHT_1,
        brightness=0.0, bold=False, align=PP_ALIGN.LEFT):
    """Add a positioned text box."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    if theme_color is not None:
        _apply_theme_color(p.font.color, theme_color, brightness)
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return box


def rect(slide, l, t, w, h, fill_theme, fill_brightness=0.0,
         text="", sz=14, fc_theme=MSO_THEME_COLOR.LIGHT_1, fc_brightness=0.0,
         bold=False):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    _apply_theme_color(shape.fill.fore_color, fill_theme, fill_brightness)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        if fc_theme is not None:
            _apply_theme_color(p.font.color, fc_theme, fc_brightness)
        p.font.bold = bold
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def header(slide, title, subtitle=""):
    """Set the slide title via the layout's title placeholder.
    Falls back to a floating text box if no title placeholder exists.
    Subtitle is always a floating text box positioned below the title.
    """
    set_bg(slide)
    # Try to use the layout's title placeholder (idx 0)
    title_set = False
    if slide.placeholders:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # Title placeholder
                ph.text = title
                title_set = True
                break
    if not title_set:
        txt(slide, 0.5, 0.3, 12, 0.6, title, sz=28, bold=True)
    if subtitle:
        txt(slide, 0.5, 0.85, 12, 0.4, subtitle, sz=14)


def card(slide, l, t, w, h, title, lines, accent_theme, tsz=16, bsz=14):
    """Add an info card with accent top border and bullet lines."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    bg.fill.background()
    bg.line.width = Pt(1.5)
    _apply_theme_color(bg.line.color, accent_theme)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(l + 0.05), Inches(t + 0.05),
                                 Inches(w - 0.1), Inches(0.06))
    bar.fill.solid()
    _apply_theme_color(bar.fill.fore_color, accent_theme)
    bar.line.fill.background()
    txt(slide, l + 0.15, t + 0.18, w - 0.3, 0.35, title,
        sz=tsz, theme_color=accent_theme, bold=True)
    body_box = slide.shapes.add_textbox(
        Inches(l + 0.15), Inches(t + 0.55),
        Inches(w - 0.3), Inches(h - 0.7))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(bsz)
        _apply_theme_color(p.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)


def bullet_list(slide, l, t, w, h, items, sz=14):
    """Add a bullet list with proper formatting."""
    body_box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = body_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(sz)
        _apply_theme_color(p.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)


# ── Plan Parsing ──

def parse_presentation_plan(file_path):
    """Parse a markdown presentation plan and extract slide data."""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    # Extract title
    title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
    presentation_title = title_match.group(1) if title_match else "Presentation"
    
    # Extract slides with flexible pattern matching
    slides = []
    
    # Pattern to match slide headers and content
    slide_pattern = r'### (?:Slide|Diapositive)\s*\d*:?\s*(.+?)(?:\n|$)(.*?)(?=\n###|\Z)'
    
    for match in re.finditer(slide_pattern, content, re.DOTALL | re.MULTILINE):
        slide_title = match.group(1).strip()
        slide_content = match.group(2).strip()
        
        # Extract layout if specified
        layout = "title-only"  # default
        layout_match = re.search(r'\*\*(?:Layout|Mise en Page)\*\*:\s*(.+?)(?:\n|$)', slide_content, re.IGNORECASE)
        if layout_match:
            layout = layout_match.group(1).strip()
            # Remove layout line from content
            slide_content = re.sub(r'\*\*(?:Layout|Mise en Page)\*\*:\s*.+?(?:\n|$)', '', slide_content, flags=re.IGNORECASE).strip()
        
        # Extract content sections
        content_text = ""
        notes_text = ""
        
        # Extract main content
        content_match = re.search(r'\*\*(?:Content|Contenu)\*\*:\s*(.*?)(?=\n\*\*|\n###|\Z)', slide_content, re.DOTALL | re.IGNORECASE)
        if content_match:
            content_text = content_match.group(1).strip()
        
        # Extract notes if present
        notes_match = re.search(r'\*\*(?:Notes|Remarques)\*\*:\s*(.*?)(?=\n\*\*|\n###|\Z)', slide_content, re.DOTALL | re.IGNORECASE)
        if notes_match:
            notes_text = notes_match.group(1).strip()
        
        # If no structured content found, treat entire slide content as content
        if not content_text and slide_content:
            content_text = slide_content
        
        slides.append({
            'title': slide_title,
            'layout': layout.lower(),
            'content': content_text,
            'notes': notes_text
        })
    
    return presentation_title, slides


def create_dynamic_slide(slide, slide_data):
    """Create a visually impressive slide with advanced design elements."""
    title = slide_data['title']
    content = slide_data['content']
    layout = slide_data['layout']
    
    # Always set the title with enhanced styling
    header(slide, title)
    
    if not content:
        # Add visual element to title-only slides
        add_title_visual_elements(slide, title)
        return
    
    # Analyze content type and choose appropriate layout
    lines = [line.strip() for line in content.split('\n') if line.strip()]
    
    # Check for different content patterns with enhanced visuals
    if lines and (lines[0].startswith('-') or lines[0].startswith('*') or lines[0].startswith('•')):
        # Enhanced bullet list with icons and colors
        bullets = [line.lstrip('-*• ').strip() for line in lines if line.strip().startswith(('-', '*', '•'))]
        if bullets:
            create_enhanced_bullet_list(slide, bullets)
    
    elif len(lines) <= 3 and all(len(line) < 100 for line in lines):
        # Enhanced centered text with visual elements
        create_centered_spotlight(slide, lines, title)
    
    elif ':' in content and len(lines) > 2:
        # Enhanced feature cards with icons and gradients
        create_feature_showcase(slide, lines)
    
    elif any(keyword in content.lower() for keyword in ['benefit', 'advantage', 'improve', 'gain']):
        # Benefits showcase with icons
        create_benefits_showcase(slide, lines)
    
    elif any(keyword in content.lower() for keyword in ['step', 'process', 'workflow', 'phase']):
        # Process flow visualization
        create_process_flow(slide, lines)
    
    else:
        # Enhanced content layout with visual hierarchy
        create_enhanced_content_layout(slide, lines, title)
    
    # Add slide number and visual footer
    add_slide_footer(slide)


def add_title_visual_elements(slide, title):
    """Add visual elements to title-only slides."""
    # Add decorative shapes based on title content
    if any(keyword in title.lower() for keyword in ['introduction', 'welcome', 'overview']):
        # Add welcoming gradient shapes
        create_gradient_background(slide, ACCENT_PRIMARY, ACCENT_INFO)
    elif any(keyword in title.lower() for keyword in ['conclusion', 'thank', 'summary']):
        # Add closing visual elements
        create_closing_visuals(slide)
    else:
        # Add standard visual accents
        add_visual_accents(slide)


def create_gradient_background(slide, color1, color2):
    """Create a subtle gradient background effect."""
    # Add overlapping gradient rectangles
    gradient1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6), Inches(13.33), Inches(1.5))
    gradient1.fill.solid()
    _apply_theme_color(gradient1.fill.fore_color, color1, -0.8)
    gradient1.line.fill.background()
    
    gradient2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), Inches(0), Inches(2.33), Inches(7))
    gradient2.fill.solid()
    _apply_theme_color(gradient2.fill.fore_color, color2, -0.9)
    gradient2.line.fill.background()


def create_closing_visuals(slide):
    """Add visual elements for closing slides."""
    # Add thank you visual elements
    for i in range(3):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                       Inches(2 + i * 3), Inches(5.5), 
                                       Inches(0.3), Inches(0.3))
        circle.fill.solid()
        _apply_theme_color(circle.fill.fore_color, [ACCENT_POSITIVE, ACCENT_INFO, ACCENT_PRIMARY][i])
        circle.line.fill.background()


def add_visual_accents(slide):
    """Add visual accent elements."""
    # Add corner accents
    accent_colors = [ACCENT_PRIMARY, ACCENT_INFO, ACCENT_POSITIVE]
    positions = [(0, 0), (11.5, 0), (0, 6.5)]
    
    for i, (x, y) in enumerate(positions):
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                       Inches(x), Inches(y), 
                                       Inches(0.5), Inches(0.5))
        accent.fill.solid()
        _apply_theme_color(accent.fill.fore_color, accent_colors[i], -0.7)
        accent.line.fill.background()


def create_enhanced_bullet_list(slide, bullets):
    """Create visually enhanced bullet points with icons and colors."""
    # Add section header
    if len(bullets) > 0:
        section_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(0.5), Inches(1.2), 
                                            Inches(12.3), Inches(0.4))
        section_rect.fill.solid()
        _apply_theme_color(section_rect.fill.fore_color, ACCENT_PRIMARY, -0.8)
        section_rect.line.fill.background()
        
        txt(slide, 0.7, 1.25, 12, 0.3, "Key Points", sz=14, 
            theme_color=ACCENT_POSITIVE, bold=True)
    
    # Create enhanced bullet points
    colors = [ACCENT_PRIMARY, ACCENT_INFO, ACCENT_POSITIVE, ACCENT_WARNING, ACCENT_AI]
    icons = ["●", "■", "▲", "◆", "★"]
    
    for i, bullet in enumerate(bullets[:6]):  # Max 6 bullets
        y_pos = 1.8 + i * 0.7
        
        # Add icon circle
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(0.7), Inches(y_pos), 
                                        Inches(0.3), Inches(0.3))
        icon_bg.fill.solid()
        _apply_theme_color(icon_bg.fill.fore_color, colors[i % len(colors)])
        icon_bg.line.fill.background()
        
        # Add icon text
        txt(slide, 0.75, y_pos + 0.05, 0.2, 0.2, icons[i % len(icons)], 
            sz=12, theme_color=MSO_THEME_COLOR.LIGHT_1, bold=True, align=PP_ALIGN.CENTER)
        
        # Add bullet text with background
        text_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(1.2), Inches(y_pos - 0.05), 
                                        Inches(11), Inches(0.4))
        text_bg.fill.solid()
        _apply_theme_color(text_bg.fill.fore_color, MSO_THEME_COLOR.LIGHT_2, -0.95)
        text_bg.line.fill.background()
        
        # Add bullet text
        txt(slide, 1.4, y_pos, 10.5, 0.3, bullet, sz=13)


def create_centered_spotlight(slide, lines, title):
    """Create a centered spotlight effect for key messages."""
    # Add spotlight background
    spotlight = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(3), Inches(2), 
                                     Inches(7.33), Inches(3))
    spotlight.fill.solid()
    _apply_theme_color(spotlight.fill.fore_color, ACCENT_PRIMARY, -0.9)
    spotlight.line.fill.background()
    
    # Add decorative elements
    for i in range(4):
        decoration = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(1 + i * 3), Inches(1.5 + i * 0.3), 
                                          Inches(0.1), Inches(4))
        decoration.fill.solid()
        _apply_theme_color(decoration.fill.fore_color, ACCENT_INFO, -0.8)
        decoration.line.fill.background()
    
    # Add centered text with emphasis
    for i, line in enumerate(lines):
        y_pos = 3.0 + i * 0.8
        txt(slide, 0.5, y_pos, 12.3, 0.6, line, sz=18, 
            theme_color=ACCENT_POSITIVE, bold=True, align=PP_ALIGN.CENTER)


def create_feature_showcase(slide, lines):
    """Create an impressive feature showcase with cards."""
    # Parse features from key-value pairs
    features = []
    for line in lines:
        if ':' in line:
            feature, description = line.split(':', 1)
            features.append((feature.strip(), description.strip()))
    
    # Create feature cards with icons
    card_width = 3.8
    card_height = 2.2
    colors = [ACCENT_PRIMARY, ACCENT_INFO, ACCENT_POSITIVE, ACCENT_WARNING]
    
    for i, (feature, description) in enumerate(features[:4]):
        x_pos = 0.7 + (i % 2) * 4.5
        y_pos = 1.5 + (i // 2) * 2.5
        
        # Create card with gradient effect
        card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x_pos), Inches(y_pos), 
                                        Inches(card_width), Inches(card_height))
        card_bg.fill.solid()
        _apply_theme_color(card_bg.fill.fore_color, colors[i % len(colors)], -0.85)
        card_bg.line.width = Pt(2)
        _apply_theme_color(card_bg.line.color, colors[i % len(colors)])
        
        # Add icon header
        icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                            Inches(x_pos + 0.2), Inches(y_pos + 0.1), 
                                            Inches(0.4), Inches(0.4))
        icon_circle.fill.solid()
        _apply_theme_color(icon_circle.fill.fore_color, colors[i % len(colors)])
        icon_circle.line.fill.background()
        
        # Add feature title
        txt(slide, x_pos + 0.7, y_pos + 0.15, card_width - 1, 0.4, 
            feature, sz=14, theme_color=colors[i % len(colors)], bold=True)
        
        # Add description
        txt(slide, x_pos + 0.2, y_pos + 0.6, card_width - 0.4, 1.2, 
            description, sz=11, theme_color=MSO_THEME_COLOR.LIGHT_1, brightness=-0.1)


def create_benefits_showcase(slide, lines):
    """Create a benefits showcase with checkmarks and progress bars."""
    # Add header
    header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.5), Inches(1.2), 
                                        Inches(12.3), Inches(0.5))
    header_rect.fill.solid()
    _apply_theme_color(header_rect.fill.fore_color, ACCENT_POSITIVE, -0.8)
    header_rect.line.fill.background()
    
    txt(slide, 0.7, 1.3, 12, 0.3, "Key Benefits & Advantages", sz=16, 
        theme_color=MSO_THEME_COLOR.LIGHT_1, bold=True)
    
    # Create benefit items with progress indicators
    for i, line in enumerate(lines[:5]):
        y_pos = 2.0 + i * 0.8
        
        # Checkmark circle
        check_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                             Inches(0.7), Inches(y_pos), 
                                             Inches(0.3), Inches(0.3))
        check_circle.fill.solid()
        _apply_theme_color(check_circle.fill.fore_color, ACCENT_POSITIVE)
        check_circle.line.fill.background()
        
        # Checkmark
        txt(slide, 0.75, y_pos + 0.05, 0.2, 0.2, "✓", 
            sz=14, theme_color=MSO_THEME_COLOR.LIGHT_1, bold=True, align=PP_ALIGN.CENTER)
        
        # Progress bar background
        progress_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(1.2), Inches(y_pos + 0.1), 
                                            Inches(8), Inches(0.2))
        progress_bg.fill.solid()
        _apply_theme_color(progress_bg.fill.fore_color, MSO_THEME_COLOR.LIGHT_2, -0.9)
        progress_bg.line.fill.background()
        
        # Progress bar fill
        progress_fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             Inches(1.2), Inches(y_pos + 0.1), 
                                             Inches(6 + i * 0.5), Inches(0.2))
        progress_fill.fill.solid()
        _apply_theme_color(progress_fill.fill.fore_color, ACCENT_POSITIVE, -0.3)
        progress_fill.line.fill.background()
        
        # Benefit text
        txt(slide, 1.4, y_pos - 0.2, 10, 0.3, line, sz=13)


def create_process_flow(slide, lines):
    """Create a visual process flow diagram."""
    # Extract steps from content
    steps = []
    for line in lines:
        if ':' in line:
            step, desc = line.split(':', 1)
            steps.append((step.strip(), desc.strip()))
        elif line.strip():
            steps.append((f"Step {len(steps) + 1}", line.strip()))
    
    # Create process flow
    box_width = 2.5
    box_height = 1.2
    arrow_length = 1.0
    
    for i, (step, desc) in enumerate(steps[:4]):
        x_pos = 1.5 + i * (box_width + arrow_length)
        y_pos = 2.5
        
        # Process box
        process_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(x_pos), Inches(y_pos), 
                                            Inches(box_width), Inches(box_height))
        process_box.fill.solid()
        _apply_theme_color(process_box.fill.fore_color, ACCENT_INFO, -0.8)
        process_box.line.width = Pt(2)
        _apply_theme_color(process_box.line.color, ACCENT_INFO)
        
        # Step number
        number_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                              Inches(x_pos + 0.1), Inches(y_pos + 0.1), 
                                              Inches(0.4), Inches(0.4))
        number_circle.fill.solid()
        _apply_theme_color(number_circle.fill.fore_color, ACCENT_PRIMARY)
        number_circle.line.fill.background()
        
        txt(slide, x_pos + 0.2, y_pos + 0.15, 0.2, 0.2, str(i + 1), 
            sz=12, theme_color=MSO_THEME_COLOR.LIGHT_1, bold=True, align=PP_ALIGN.CENTER)
        
        # Step text
        txt(slide, x_pos + 0.6, y_pos + 0.3, box_width - 0.8, 0.3, 
            step, sz=12, theme_color=ACCENT_PRIMARY, bold=True)
        
        # Description
        txt(slide, x_pos + 0.2, y_pos + 0.7, box_width - 0.4, 0.4, 
            desc, sz=10)
        
        # Arrow to next step
        if i < len(steps) - 1 and i < 3:
            arrow_x = x_pos + box_width
            txt(slide, arrow_x, y_pos + 0.4, arrow_length, 0.3, "→", 
                sz=24, theme_color=ACCENT_POSITIVE, bold=True, align=PP_ALIGN.CENTER)


def create_enhanced_content_layout(slide, lines, title):
    """Create an enhanced content layout with visual hierarchy."""
    # Add title accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.5), Inches(1.2), 
                                       Inches(3), Inches(0.1))
    accent_bar.fill.solid()
    _apply_theme_color(accent_bar.fill.fore_color, ACCENT_PRIMARY)
    accent_bar.line.fill.background()
    
    # Group content into sections
    sections = []
    current_section = []
    
    for line in lines:
        if len(line) < 50 and not line.endswith(':'):
            # Likely a header
            if current_section:
                sections.append(current_section)
            current_section = [line]
        else:
            current_section.append(line)
    
    if current_section:
        sections.append(current_section)
    
    # Create visual sections
    colors = [ACCENT_PRIMARY, ACCENT_INFO, ACCENT_POSITIVE]
    
    for i, section in enumerate(sections[:3]):
        y_pos = 1.5 + i * 2.0
        
        # Section background
        section_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(0.7), Inches(y_pos), 
                                           Inches(11.9), Inches(1.6))
        section_bg.fill.solid()
        _apply_theme_color(section_bg.fill.fore_color, colors[i % len(colors)], -0.9)
        section_bg.line.width = Pt(1)
        _apply_theme_color(section_bg.line.color, colors[i % len(colors)], -0.5)
        
        # Section content
        for j, line in enumerate(section[:2]):
            text_y = y_pos + 0.2 + j * 0.6
            if j == 0 and len(line) < 50:
                # Header
                txt(slide, 1.0, text_y, 11, 0.4, line, 
                    sz=14, theme_color=colors[i % len(colors)], bold=True)
            else:
                # Content
                txt(slide, 1.0, text_y, 11, 0.5, line, sz=12)


def add_slide_footer(slide):
    """Add a professional slide footer with slide number."""
    # Add footer line
    footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.5), Inches(6.8), 
                                        Inches(12.3), Inches(0.05))
    footer_line.fill.solid()
    _apply_theme_color(footer_line.fill.fore_color, ACCENT_INFO, -0.7)
    footer_line.line.fill.background()
    
    # Add decorative element
    decor = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(12), Inches(6.7), 
                                  Inches(0.3), Inches(0.3))
    decor.fill.solid()
    _apply_theme_color(decor.fill.fore_color, ACCENT_POSITIVE)
    decor.line.fill.background()


def generate_presentation_from_plan(plan_file_path, output_dir):
    """Generate a dynamic PowerPoint presentation from a markdown plan."""
    
    # Parse the plan
    presentation_title, slides = parse_presentation_plan(plan_file_path)
    
    # Create presentation with template
    template_path = Path(__file__).resolve().parent / "template.pptx"
    if template_path.exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
    
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Create slides dynamically
    for i, slide_data in enumerate(slides):
        if i == 0:
            # First slide - check if it should be a title slide
            if slide_data['layout'] in ['title', 'title slide', 'cover']:
                slide = prs.slides.add_slide(get_cover_layout(prs))
                # Populate cover slide placeholders
                for ph in slide.placeholders:
                    idx = ph.placeholder_format.idx
                    if idx == 0:  # Title
                        ph.text = slide_data['title']
                    elif idx == 1 and slide_data['content']:  # Subtitle
                        # Extract first line as subtitle
                        first_line = slide_data['content'].split('\n')[0].strip()
                        ph.text = first_line
            else:
                slide = prs.slides.add_slide(get_title_only_layout(prs))
                create_dynamic_slide(slide, slide_data)
        else:
            # Content slides
            slide = prs.slides.add_slide(get_title_only_layout(prs))
            create_dynamic_slide(slide, slide_data)
    
    # Add closing slide if not already present
    if not any('thank' in slide['title'].lower() or 'conclusion' in slide['title'].lower() 
               for slide in slides):
        slide = prs.slides.add_slide(get_closing_layout(prs))
        header(slide, "Thank You", f"{presentation_title}")
    
    # Set document properties
    prs.core_properties.title = presentation_title
    prs.core_properties.subject = f"Generated from plan: {Path(plan_file_path).name}"
    prs.core_properties.comments = "Generated by generate-pptx-from-plan skill"
    
    # Generate output filename
    safe_title = re.sub(r'[^\w\s-]', '', presentation_title).strip()
    safe_title = re.sub(r'[-\s]+', '_', safe_title)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    output_filename = f"{safe_title}_{timestamp}.pptx"
    
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, output_filename)
    
    # Save the presentation
    prs.save(output_path)
    return output_path


def main():
    """Main function for command line usage."""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python generate_pptx.py <plan_file_path> [output_directory]")
        print("Example: python generate_pptx.py presentation-plan.md ./output")
        return
    
    plan_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."
    
    if not os.path.exists(plan_file):
        print(f"Error: Plan file '{plan_file}' not found.")
        return
    
    try:
        output_path = generate_presentation_from_plan(plan_file, output_dir)
        print(f"✅ Presentation generated: {output_path}")
    except Exception as e:
        print(f"❌ Error generating presentation: {str(e)}")


if __name__ == "__main__":
    main()
