#!/usr/bin/env python3
"""
Template: Comparison PPTX Generator
====================================
Reference script for the generate-comparison-pptx skill.
Generates a professional PowerPoint presentation comparing multiple items
across dimensions with native editable charts.

THEME-AWARE: All colors use PowerPoint theme slots (Accent 1-6, Text 1/2,
Background 1/2) so the user can switch the theme in PowerPoint and all
colors adapt automatically. No hardcoded RGB values in slide content.

Usage:
  1. Copy this script to your working directory
  2. Replace the DATA SECTION with your actual comparison data
  3. Run: python generate_pptx.py
  4. Open the PPTX → Design tab → pick any theme to restyle

Requirements: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pathlib import Path


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  THEME COLOR MAPPING                                              ║
# ║                                                                   ║
# ║  All colors reference PowerPoint theme slots. When the user       ║
# ║  switches the theme (Design tab), every element recolors.         ║
# ║                                                                   ║
# ║  Semantic role        → Theme slot        → Default appearance    ║
# ║  ─────────────────────────────────────────────────────────────     ║
# ║  Slide background     → (inherited)       → from theme            ║
# ║  Card background      → LIGHT_2 +0.4 bright → subtle mid-tone     ║
# ║  Primary text         → LIGHT_1           → white                ║
# ║  Body text            → LIGHT_1 -0.25 bright → light gray        ║
# ║  Subtitle/muted text  → LIGHT_2           → mid gray             ║
# ║  Item colors          → ACCENT_1..6       → per-item accent      ║
# ║  Insight box bg       → ACCENT_1 -0.7 bright → dark tinted       ║
# ╚═══════════════════════════════════════════════════════════════════╝

# Theme accent slots for items (up to 6 natively, extend with brightness)
ACCENT_SLOTS = [
    MSO_THEME_COLOR.ACCENT_1,  # Blue
    MSO_THEME_COLOR.ACCENT_2,  # Green (or theme's 2nd accent)
    MSO_THEME_COLOR.ACCENT_3,  # Purple
    MSO_THEME_COLOR.ACCENT_4,  # Orange
    MSO_THEME_COLOR.ACCENT_5,  # Red
    MSO_THEME_COLOR.ACCENT_6,  # Teal
]

# For items 7-8, reuse accents with brightness shift
ACCENT_SLOTS_EXTENDED = ACCENT_SLOTS + [
    (MSO_THEME_COLOR.ACCENT_1, 0.4),   # Lighter variant of Accent 1
    (MSO_THEME_COLOR.ACCENT_2, 0.4),   # Lighter variant of Accent 2
]

# XML scheme color names matching MSO_THEME_COLOR for chart coloring
SCHEME_CLR_NAMES = ["accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  DATA SECTION — replace everything below with your actual data   ║
# ╚═══════════════════════════════════════════════════════════════════╝

# Title slide content
TITLE = "Item Comparison"
SUBTITLE = "N Items  •  Same Conditions  •  Real Results"
FOOTER = "Your context line  |  Date"
DATA_SOURCE = "Data: describe where the data comes from"
AI_NOTICE = "AI-generated presentation — data verified by author"

# Items to compare (2–8 items)
ITEMS = ["Item A", "Item B", "Item C", "Item D", "Item E"]

# --- Slide 1: Two single-metric bar charts side by side ---
METRIC_1_LABEL = "Cost ($)"
METRIC_1_VALUES = [2.90, 14.23, 0.99, 0.03, 0.03]
METRIC_1_FORMAT = '$#,##0.00'
METRIC_1_AXIS = "Cost (USD)"

METRIC_2_LABEL = "Duration (min)"
METRIC_2_VALUES = [2.7, 3.5, 1.8, 0.7, 0.7]
METRIC_2_FORMAT = '#,##0.0 "min"'
METRIC_2_AXIS = "Duration (minutes)"

SLIDE_1_INSIGHT = "💡 Item C = best value  |  Item B = deepest analysis"

# --- Slide 2: Multi-dimension grouped bar chart ---
DIMENSIONS = ["Dim 1", "Dim 2", "Dim 3", "Dim 4", "Dim 5", "Dim 6", "Overall"]
SCORES = {
    "Item A": [7, 5, 5, 6, 7, 7, 6],
    "Item B": [5, 3, 4, 5, 5, 5, 4],
    "Item C": [7, 5, 4, 5, 5, 7, 6],
    "Item D": [6, 5, 4, 5, 6, 7, 6],
    "Item E": [6, 7, 6, 5, 8, 7, 7],
}
SCORE_MAX = 10
SLIDE_2_INSIGHTS = [
    ("⚠️", "Item E scores highest but may be inflated"),
    ("✅", "Item B scores lowest = most thorough analysis"),
]

# --- Slide 3: Found/missed stacked bar + key findings ---
DETECTION_ORDER = ["Item B", "Item C", "Item A", "Item D", "Item E"]
FOUND_COUNTS    = [15, 11, 6, 4, 0]
TOTAL_POSSIBLE  = 15

KEY_FINDINGS = [
    ("CRITICAL", "Description of critical finding", MSO_THEME_COLOR.ACCENT_5),
    ("HIGH",     "Description of high-severity finding", MSO_THEME_COLOR.ACCENT_4),
    ("HIGH",     "Another high-severity finding", MSO_THEME_COLOR.ACCENT_4),
    ("MEDIUM",   "Description of medium finding", MSO_THEME_COLOR.ACCENT_1),
    ("SECURITY", "Description of security finding", MSO_THEME_COLOR.ACCENT_3),
]
SLIDE_3_INSIGHT = "Only Item B found ALL issues  •  Item C found 11 at 7% cost"

# --- Slide 4: Depth metric bar chart + profile cards ---
DEPTH_LABEL = "Input Tokens (K)"
DEPTH_VALUES = [913, 882, 937, 37, 25]
DEPTH_FORMAT = '#,##0"K"'
DEPTH_AXIS = "Input Tokens (thousands)"

PROFILES = [
    ("Item B",     "882K tokens • 50+ tool calls\nFinds every issue\nBest for: deep audits", 1),
    ("Item C ⭐",  "937K tokens • 43+ tool calls\n73% of B's findings at 7% cost\nBest for: routine use", 2),
    ("Item A",     "913K tokens • ~17 tool calls\nGood balance\nBest for: mid-tier needs", 0),
    ("Item D",     "37K tokens • few tool calls\nFast triage\nBest for: smoke tests", 3),
    ("Item E",     "25K tokens • minimal reads\nNot recommended", 4),
]

# --- Slide 5: Recommendation tiers ---
RECOMMENDATIONS = [
    {
        "title": "🏆  DEFAULT — Item C",
        "subtitle": "$0.99  •  1.8 min  •  11/15 issues",
        "bullets": [
            "Best ROI: 73% of top findings at 7% cost",
            "Catches critical issues reliably",
            "Use for: all routine work",
        ],
        "accent_idx": 2,  # index into ACCENT_SLOTS
    },
    {
        "title": "🔬  DEEP — Item B",
        "subtitle": "$14.23  •  3.5 min  •  15/15 issues",
        "bullets": [
            "Finds everything with evidence",
            "Most specific and thorough",
            "Use for: audits, compliance, critical work",
        ],
        "accent_idx": 1,
    },
    {
        "title": "⚡  QUICK — Item D",
        "subtitle": "$0.03  •  0.7 min  •  4/15 issues",
        "bullets": [
            "Ultra-cheap bulk scanning",
            "Catches only the biggest issues",
            "Use for: triage, smoke tests",
        ],
        "accent_idx": 3,
    },
]
RECOMMENDATION_FOOTER = "Item A sits between C and B — use when C isn't enough  |  Item E: skip"

# Output path
# Sources: list of (label, url) — extracted from source data/docs
# A "Sources & References" slide will be auto-generated at the end.
SOURCES = [
    # ("Label shown on slide", "https://full-url"),
]

# Template: path to a .pptx template file (optional).
# If provided, the presentation inherits the template's theme and branding.
# If None or file not found, falls back to a blank presentation.
TEMPLATE_PATH = Path(__file__).parent / "template.pptx"

OUTPUT_PATH = Path(__file__).parent / "comparison.pptx"


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  HELPER FUNCTIONS                                                 ║
# ╚═══════════════════════════════════════════════════════════════════╝

def _apply_theme_color(color_format, theme_color, brightness=0.0):
    """Apply a theme color with optional brightness adjustment."""
    color_format.theme_color = theme_color
    if brightness != 0.0:
        color_format.brightness = brightness


def _get_item_accent(index):
    """Get the theme color slot for item at given index.
    Returns (theme_color, brightness) tuple.
    """
    entry = ACCENT_SLOTS_EXTENDED[index % len(ACCENT_SLOTS_EXTENDED)]
    if isinstance(entry, tuple):
        return entry
    return (entry, 0.0)


def set_slide_bg(slide):
    """No-op: let the slide inherit its background from the slide master/layout.
    This ensures that when the user switches the PowerPoint theme, the
    background changes along with everything else. If you explicitly set
    a background here, it becomes an override that blocks theme inheritance.
    """
    pass


def add_text_box(slide, left, top, width, height, text,
                 font_size=14, theme_color=MSO_THEME_COLOR.LIGHT_1,
                 brightness=0.0, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a positioned text box. Default color is LIGHT_1 (white on dark themes).
    For textboxes (not placeholders), PowerPoint does NOT inherit the
    master's text color, so we must set it explicitly."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if theme_color is not None:
        _apply_theme_color(p.font.color, theme_color, brightness)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height,
                     fill_theme, fill_brightness=0.0,
                     text="", font_size=12,
                     fc_theme=MSO_THEME_COLOR.LIGHT_1, fc_brightness=0.0,
                     bold=False):
    """Add a rounded rectangle. Text defaults to LIGHT_1 (white)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    _apply_theme_color(shape.fill.fore_color, fill_theme, fill_brightness)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        if fc_theme is not None:
            _apply_theme_color(p.font.color, fc_theme, fc_brightness)
        p.font.bold = bold
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_card(slide, left, top, width, height, title, lines,
             accent_theme, title_size=14, body_size=11):
    """Add an info card with accent top border and bullet lines.
    All body lines are rendered as paragraphs inside a single text frame.
    Card background is transparent — inherits from slide/theme background.
    A subtle border provides visual separation on any theme.
    """
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.background()  # transparent
    card.line.width = Pt(1.5)
    _apply_theme_color(card.line.color, accent_theme)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + 0.05), Inches(top + 0.05),
        Inches(width - 0.1), Inches(0.06)
    )
    bar.fill.solid()
    _apply_theme_color(bar.fill.fore_color, accent_theme)
    bar.line.fill.background()
    add_text_box(slide, left + 0.15, top + 0.15, width - 0.3, 0.35,
                 title, font_size=title_size, theme_color=accent_theme, bold=True)
    body_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.55),
        Inches(width - 0.3), Inches(height - 0.7)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(body_size)
        _apply_theme_color(p.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)
    return card


def _get_solid_fill_element(fill_obj):
    """Get the XML solidFill element, handling different python-pptx internals.
    Some versions return an lxml element directly from fill._fill, while
    others wrap it in a _SolidFill object with a _solidFill attribute.
    """
    sf = fill_obj._fill
    if hasattr(sf, 'tag'):
        return sf
    if hasattr(sf, '_solidFill'):
        return sf._solidFill
    if hasattr(sf, '_element'):
        return sf._element
    return sf


def _set_chart_point_theme_color(point, scheme_clr_name, brightness=0.0):
    """Set a chart data point's fill to a theme color via XML.
    python-pptx charts don't natively support theme colors on data points,
    so we manipulate the XML directly to use schemeClr instead of srgbClr.
    """
    fill = point.format.fill
    fill.solid()
    solid_fill = _get_solid_fill_element(fill)
    # Remove any existing color child
    for child in list(solid_fill):
        if child.tag.endswith(('srgbClr', 'schemeClr', 'sysClr')):
            solid_fill.remove(child)
    # Add schemeClr element
    scheme_clr = solid_fill.makeelement(qn('a:schemeClr'), {'val': scheme_clr_name})
    if brightness != 0.0:
        if brightness > 0:
            mod = scheme_clr.makeelement(qn('a:lumMod'),
                                          {'val': str(int((1 - brightness) * 100000))})
            off = scheme_clr.makeelement(qn('a:lumOff'),
                                          {'val': str(int(brightness * 100000))})
            scheme_clr.append(mod)
            scheme_clr.append(off)
        else:
            mod = scheme_clr.makeelement(qn('a:lumMod'),
                                          {'val': str(int((1 + brightness) * 100000))})
            scheme_clr.append(mod)
    solid_fill.append(scheme_clr)


def _set_series_theme_color(series, scheme_clr_name, brightness=0.0):
    """Set an entire chart series fill to a theme color via XML."""
    fill = series.format.fill
    fill.solid()
    solid_fill = _get_solid_fill_element(fill)
    for child in list(solid_fill):
        if child.tag.endswith(('srgbClr', 'schemeClr', 'sysClr')):
            solid_fill.remove(child)
    scheme_clr = solid_fill.makeelement(qn('a:schemeClr'), {'val': scheme_clr_name})
    if brightness != 0.0:
        if brightness > 0:
            mod = scheme_clr.makeelement(qn('a:lumMod'),
                                          {'val': str(int((1 - brightness) * 100000))})
            off = scheme_clr.makeelement(qn('a:lumOff'),
                                          {'val': str(int(brightness * 100000))})
            scheme_clr.append(mod)
            scheme_clr.append(off)
        else:
            mod = scheme_clr.makeelement(qn('a:lumMod'),
                                          {'val': str(int((1 + brightness) * 100000))})
            scheme_clr.append(mod)
    solid_fill.append(scheme_clr)


def _get_scheme_clr_name(item_index):
    """Get the XML scheme color name for an item index."""
    if item_index < len(SCHEME_CLR_NAMES):
        return SCHEME_CLR_NAMES[item_index]
    # For items beyond 6, reuse with brightness (handled by caller)
    return SCHEME_CLR_NAMES[item_index % len(SCHEME_CLR_NAMES)]


def _get_item_brightness(item_index):
    """Get brightness adjustment for items beyond the 6 native accents."""
    if item_index < len(SCHEME_CLR_NAMES):
        return 0.0
    return 0.4  # lighter variant for overflow items


def style_bar_chart(chart, num_items, data_label_format='#,##0',
                    axis_label=None, max_scale=None):
    """Apply consistent theme-aware styling to a single-series bar chart."""
    chart.has_legend = False
    chart.style = 2
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    for i in range(num_items):
        pt = series.points[i]
        _set_chart_point_theme_color(
            pt, _get_scheme_clr_name(i), _get_item_brightness(i))
    series.data_labels.show_value = True
    series.data_labels.font.size = Pt(12)
    _apply_theme_color(series.data_labels.font.color, MSO_THEME_COLOR.LIGHT_1)
    series.data_labels.font.bold = True
    series.data_labels.number_format = data_label_format
    chart.category_axis.tick_labels.font.size = Pt(10)
    _apply_theme_color(chart.category_axis.tick_labels.font.color,
                       MSO_THEME_COLOR.LIGHT_1, -0.25)
    chart.value_axis.tick_labels.font.size = Pt(9)
    _apply_theme_color(chart.value_axis.tick_labels.font.color,
                       MSO_THEME_COLOR.LIGHT_1, -0.4)
    if axis_label:
        chart.value_axis.has_title = True
        ax_title = chart.value_axis.axis_title.text_frame.paragraphs[0]
        ax_title.text = axis_label
        ax_title.font.size = Pt(10)
        _apply_theme_color(ax_title.font.color, MSO_THEME_COLOR.LIGHT_1, -0.4)
    if max_scale is not None:
        chart.value_axis.maximum_scale = max_scale


def add_slide_header(slide, title, subtitle=""):
    """Add standard title + subtitle to a slide."""
    set_slide_bg(slide)
    add_text_box(slide, 0.5, 0.3, 12, 0.6, title,
                 font_size=28, bold=True)
    if subtitle:
        add_text_box(slide, 0.5, 0.85, 12, 0.4, subtitle,
                     font_size=14)


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  SLIDE BUILDERS                                                   ║
# ╚═══════════════════════════════════════════════════════════════════╝

def build_title_slide(prs):
    """Slide 0: Title with intent, data source, AI-generated note, and item badges."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 1, 1.4, 11, 1.2, TITLE,
                 font_size=40, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.8, 11, 0.6, SUBTITLE,
                 font_size=22, theme_color=MSO_THEME_COLOR.ACCENT_1,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.6, 11, 0.5, FOOTER,
                 font_size=16,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.2, 11, 0.4, DATA_SOURCE,
                 font_size=14,
                 alignment=PP_ALIGN.CENTER)
    # Item badges
    n = len(ITEMS)
    total_width = n * 1.9
    x_start = (13.333 - total_width) / 2
    for i, name in enumerate(ITEMS):
        tc, br = _get_item_accent(i)
        add_rounded_rect(slide, x_start + i * 1.9, 5.1, 1.7, 0.5,
                         tc, br, name, font_size=14, bold=True)
    # AI-generated notice
    add_text_box(slide, 1, 6.2, 11, 0.4, AI_NOTICE,
                 font_size=11,
                 alignment=PP_ALIGN.CENTER)


def build_cost_performance_slide(prs):
    """Slide 1: Two side-by-side bar charts for primary metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Cost & Performance at a Glance",
                     "Same conditions, same pipeline")

    n = len(ITEMS)

    cd1 = CategoryChartData()
    cd1.categories = ITEMS
    cd1.add_series(METRIC_1_LABEL, METRIC_1_VALUES)
    cf1 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.5), Inches(5.8), Inches(4.2), cd1
    )
    style_bar_chart(cf1.chart, n, METRIC_1_FORMAT, METRIC_1_AXIS)

    cd2 = CategoryChartData()
    cd2.categories = ITEMS
    cd2.add_series(METRIC_2_LABEL, METRIC_2_VALUES)
    cf2 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.2), cd2
    )
    style_bar_chart(cf2.chart, n, METRIC_2_FORMAT, METRIC_2_AXIS)

    add_rounded_rect(slide, 2.5, 6.0, 8.3, 0.55,
                     MSO_THEME_COLOR.ACCENT_1, -0.7,
                     SLIDE_1_INSIGHT, font_size=14,
                     fc_theme=MSO_THEME_COLOR.ACCENT_2)


def build_quality_scores_slide(prs):
    """Slide 2: Grouped bar chart for multi-dimension scoring."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Quality Scores by Dimension",
                     "Lower scores may indicate stricter analysis")

    cd = CategoryChartData()
    cd.categories = DIMENSIONS
    for item in ITEMS:
        cd.add_series(item, SCORES[item])

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.5), cd
    )
    chart = cf.chart
    chart.style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    plot.overlap = -10
    for i in range(len(ITEMS)):
        _set_series_theme_color(
            plot.series[i], _get_scheme_clr_name(i), _get_item_brightness(i))

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    _apply_theme_color(chart.legend.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
    chart.category_axis.tick_labels.font.size = Pt(10)
    _apply_theme_color(chart.category_axis.tick_labels.font.color,
                       MSO_THEME_COLOR.LIGHT_1, -0.25)
    chart.value_axis.tick_labels.font.size = Pt(9)
    _apply_theme_color(chart.value_axis.tick_labels.font.color,
                       MSO_THEME_COLOR.LIGHT_1, -0.4)
    chart.value_axis.maximum_scale = SCORE_MAX
    chart.value_axis.minimum_scale = 0

    # Insight callouts
    x = 0.8
    tints = [MSO_THEME_COLOR.ACCENT_3, MSO_THEME_COLOR.ACCENT_1]
    for idx, (emoji, text) in enumerate(SLIDE_2_INSIGHTS):
        add_rounded_rect(slide, x, 6.2, 5.5, 0.55,
                         tints[idx % len(tints)], -0.7,
                         f"{emoji} {text}", font_size=14,
                         fc_theme=MSO_THEME_COLOR.ACCENT_4)
        x += 6.0


def build_detection_slide(prs):
    """Slide 3: Stacked bar (found/missed) + key findings cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Critical Issues Detection",
                     f"{TOTAL_POSSIBLE} real issues — how many did each item find?")

    missed = [TOTAL_POSSIBLE - f for f in FOUND_COUNTS]

    cd = CategoryChartData()
    cd.categories = DETECTION_ORDER
    cd.add_series("Found", FOUND_COUNTS)
    cd.add_series("Missed", missed)

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.5), cd
    )
    chart = cf.chart
    chart.style = 2
    plot = chart.plots[0]

    s_found = plot.series[0]
    _set_series_theme_color(s_found, "accent2")  # green = found
    s_found.data_labels.show_value = True
    s_found.data_labels.font.size = Pt(11)
    _apply_theme_color(s_found.data_labels.font.color, MSO_THEME_COLOR.LIGHT_1)
    s_found.data_labels.font.bold = True

    s_missed = plot.series[1]
    _set_series_theme_color(s_missed, "dk1", 0.4)  # muted dark = missed
    s_missed.data_labels.show_value = True
    s_missed.data_labels.font.size = Pt(10)
    _apply_theme_color(s_missed.data_labels.font.color, MSO_THEME_COLOR.LIGHT_1, -0.4)

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    _apply_theme_color(chart.legend.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
    chart.category_axis.tick_labels.font.size = Pt(11)
    _apply_theme_color(chart.category_axis.tick_labels.font.color,
                       MSO_THEME_COLOR.LIGHT_1, -0.25)
    chart.value_axis.tick_labels.font.size = Pt(9)
    _apply_theme_color(chart.value_axis.tick_labels.font.color,
                       MSO_THEME_COLOR.LIGHT_1, -0.4)
    chart.value_axis.maximum_scale = TOTAL_POSSIBLE

    # Key findings cards on the right
    add_text_box(slide, 7.0, 1.3, 5.5, 0.4,
                 "Key Issues Found", font_size=18, bold=True)
    y = 1.8
    for severity, desc, accent in KEY_FINDINGS:
        add_rounded_rect(slide, 7.0, y, 1.2, 0.65, accent,
                         text=severity, font_size=12, bold=True)
        add_text_box(slide, 8.3, y + 0.05, 4.5, 0.6,
                     desc, font_size=12)
        y += 0.75

    add_text_box(slide, 0.5, 5.8, 12, 0.8, SLIDE_3_INSIGHT,
                 font_size=14, theme_color=MSO_THEME_COLOR.ACCENT_2, bold=True,
                 alignment=PP_ALIGN.CENTER)


def build_depth_slide(prs):
    """Slide 4: Depth metric bar chart + profile cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, f"Analysis Depth: {DEPTH_LABEL} vs Findings",
                     "More depth = more issues discovered")

    n = len(ITEMS)
    cd = CategoryChartData()
    cd.categories = ITEMS
    cd.add_series(DEPTH_LABEL, DEPTH_VALUES)

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.4), Inches(5.8), Inches(4.0), cd
    )
    style_bar_chart(cf.chart, n, DEPTH_FORMAT, DEPTH_AXIS)

    # Profile cards on the right — each profile's description as a single text block
    y = 1.4
    for label, desc, color_idx in PROFILES:
        tc, br = _get_item_accent(color_idx)
        add_rounded_rect(slide, 7.0, y, 1.8, 0.7, tc, br,
                         label, font_size=12, bold=True)
        # Description as single text frame with line breaks
        desc_box = slide.shapes.add_textbox(
            Inches(8.9), Inches(y + 0.02), Inches(4.0), Inches(0.7)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        desc_lines = desc.split('\n')
        for i, line in enumerate(desc_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
            _apply_theme_color(p.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
            p.font.name = "Segoe UI"
        y += 0.82


def build_recommendation_slide(prs):
    """Slide 1 (Executive Summary): Tiered recommendation cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Executive Summary: Which Item When?")

    n = len(RECOMMENDATIONS)
    card_width = 3.9
    total = n * card_width + (n - 1) * 0.2
    x = (13.333 - total) / 2

    for rec in RECOMMENDATIONS:
        tc, br = _get_item_accent(rec["accent_idx"])

        # Card outline
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.2), Inches(card_width), Inches(4.8)
        )
        card.fill.background()  # transparent
        card.line.width = Pt(2)
        _apply_theme_color(card.line.color, tc, br)

        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.05), Inches(1.25), Inches(card_width - 0.1), Inches(0.08)
        )
        bar.fill.solid()
        _apply_theme_color(bar.fill.fore_color, tc, br)
        bar.line.fill.background()

        # Title & subtitle
        add_text_box(slide, x + 0.2, 1.45, card_width - 0.4, 0.4,
                     rec["title"], font_size=18, theme_color=tc, brightness=br, bold=True)
        add_text_box(slide, x + 0.2, 1.9, card_width - 0.4, 0.35,
                     rec["subtitle"], font_size=14)

        # Bullets as a single text frame with multiple paragraphs
        bullet_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(2.4),
            Inches(card_width - 0.4), Inches(2.8)
        )
        btf = bullet_box.text_frame
        btf.word_wrap = True
        for bi, bullet in enumerate(rec["bullets"]):
            p = btf.paragraphs[0] if bi == 0 else btf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(14)
            _apply_theme_color(p.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
            p.font.name = "Segoe UI"
            p.space_after = Pt(6)

        x += card_width + 0.2

    add_rounded_rect(slide, 1.5, 6.3, 10.3, 0.55,
                     MSO_THEME_COLOR.ACCENT_1, -0.7,
                     RECOMMENDATION_FOOTER, font_size=14)


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  SOURCES SLIDE                                                    ║
# ╚═══════════════════════════════════════════════════════════════════╝

def build_sources_slide(prs):
    """Final slide: Sources & References with clickable hyperlinks."""
    if not SOURCES:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Sources & References")

    body = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5)
    )
    tf = body.text_frame
    tf.word_wrap = True
    for i, (label, url) in enumerate(SOURCES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run_bullet = p.add_run()
        run_bullet.text = "• "
        run_bullet.font.size = Pt(14)
        _apply_theme_color(run_bullet.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
        run_bullet.font.name = "Segoe UI"
        run_link = p.add_run()
        run_link.text = label
        run_link.font.size = Pt(14)
        run_link.font.name = "Segoe UI"
        _apply_theme_color(run_link.font.color, MSO_THEME_COLOR.HYPERLINK)
        run_link.font.underline = True
        run_link.hyperlink.address = url
        run_url = p.add_run()
        run_url.text = f"  ({url})"
        run_url.font.size = Pt(10)
        run_url.font.name = "Segoe UI"
        _apply_theme_color(run_url.font.color, MSO_THEME_COLOR.LIGHT_2)


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  MAIN — assemble and save                                        ║
# ╚═══════════════════════════════════════════════════════════════════╝

def main():
    # Use template if provided, otherwise blank presentation
    if TEMPLATE_PATH and Path(TEMPLATE_PATH).exists():
        prs = Presentation(str(TEMPLATE_PATH))
        print(f"📎 Using template: {TEMPLATE_PATH}")
    else:
        prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)
    build_recommendation_slide(prs)
    build_cost_performance_slide(prs)
    build_quality_scores_slide(prs)
    build_detection_slide(prs)
    build_depth_slide(prs)
    build_sources_slide(prs)

    # Set document properties
    prs.core_properties.title = TITLE
    prs.core_properties.subject = SUBTITLE
    prs.core_properties.comments = "Generated by generate-comparison-pptx skill"

    prs.save(str(OUTPUT_PATH))
    print(f"✅ Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
