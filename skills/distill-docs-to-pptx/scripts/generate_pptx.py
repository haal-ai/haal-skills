#!/usr/bin/env python3
"""
Template: Documentation Distillation PPTX Generator
=====================================================
Reference script for the distill-docs-to-pptx skill.
Generates a professional PowerPoint presentation that distills large markdown
documentation into concise, visual slides.

THEME-AWARE: All colors use PowerPoint theme slots (Accent 1-6, Text 1/2,
Background 1/2) so the user can switch the theme in PowerPoint and all
colors adapt automatically. No hardcoded RGB values in slide content.

LAYOUT-AWARE: Slide layouts are resolved by name from the template, not by
hardcoded index. The title slide uses the cover layout, content slides use
the content layout, and the closing slide uses the closing layout. Titles
are populated via the layout's built-in placeholders so they inherit the
template's font, position, and style.

Usage:
  1. Copy this script to your working directory
  2. Replace the DATA SECTION with content extracted from your docs
  3. Run: python generate_pptx.py
  4. Open the PPTX → Design tab → pick any theme to restyle

Requirements: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path


SCRIPT_DIR = Path(__file__).resolve().parent


def _presentations_root(path: Path):
    for p in [path, *path.parents]:
        if p.name == "presentations" and p.parent.name == "work" and p.parent.parent.name == ".olaf":
            return p
    return None


def _ensure_safe_execution_location():
    root = _presentations_root(SCRIPT_DIR)
    if root is None:
        raise SystemExit(
            "This file is a TEMPLATE. Do not run it from the installed skill folder. "
            "Copy it (and template.pptx) to: <repo>/.olaf/work/presentations/<presentation-name>/ "
            "and run it from there."
        )


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  THEME COLOR MAPPING                                              ║
# ║                                                                   ║
# ║  All colors reference PowerPoint theme slots. When the user       ║
# ║  switches the theme (Design tab), every element recolors.         ║
# ║                                                                   ║
# ║  Semantic role        → Theme slot        → Default appearance    ║
# ║  ─────────────────────────────────────────────────────────────     ║
# ║  Slide background     → (inherited)       → from theme            ║
# ║  Card background      → LIGHT_2 +0.4 bright → subtle mid-tone     ║
# ║  Primary text         → LIGHT_1           → white                ║
# ║  Body text            → LIGHT_1 -0.25 bright → light gray        ║
# ║  Subtitle/muted text  → LIGHT_2           → mid gray             ║
# ║  Primary accent       → ACCENT_1          → blue                 ║
# ║  Secondary accent     → ACCENT_2          → green                ║
# ║  Tertiary accent      → ACCENT_3          → purple               ║
# ║  Warning accent       → ACCENT_4          → orange               ║
# ║  Danger accent        → ACCENT_5          → red                  ║
# ║  Info accent          → ACCENT_6          → teal                 ║
# ║  Insight box bg       → ACCENT_1 -0.7 bright → dark tinted       ║
# ╚═══════════════════════════════════════════════════════════════════╝

# Semantic names for theme color slots — used in data section
ACCENT_PRIMARY   = MSO_THEME_COLOR.ACCENT_1   # Blue / structural
ACCENT_POSITIVE  = MSO_THEME_COLOR.ACCENT_2   # Green / output / quality
ACCENT_AI        = MSO_THEME_COLOR.ACCENT_3   # Purple / validation
ACCENT_WARNING   = MSO_THEME_COLOR.ACCENT_4   # Orange / config
ACCENT_DANGER    = MSO_THEME_COLOR.ACCENT_5   # Red / problems
ACCENT_INFO      = MSO_THEME_COLOR.ACCENT_6   # Teal / processing


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  LAYOUT RESOLUTION                                                ║
# ║                                                                   ║
# ║  Layouts are resolved by name from the template. This avoids      ║
# ║  hardcoded indices that break when the template changes.           ║
# ║                                                                   ║
# ║  Three layout roles:                                              ║
# ║    cover   → title/cover slide (first slide)                      ║
# ║    content → body slides with title + free content area            ║
# ║    closing → final/outro slide (sources, thank you)               ║
# ╚═══════════════════════════════════════════════════════════════════╝

# Layout name patterns for each role. The resolver tries each pattern
# (case-insensitive substring match) in order and returns the first hit.
# Customize these if your template uses different naming conventions.
LAYOUT_COVER_NAMES = ["cover slide", "cover", "title slide"]
LAYOUT_CONTENT_NAMES = ["title, subtitle and content", "title and content"]
LAYOUT_TITLE_ONLY_NAMES = ["title only", "title-only"]
LAYOUT_CLOSING_NAMES = ["closing", "end", "thank"]


def _find_layout(prs, name_patterns, fallback_index=0):
    """Find a slide layout by name pattern (case-insensitive substring match).
    Tries each pattern in order and returns the first matching layout.
    Falls back to the given index if no name matches.
    """
    for pattern in name_patterns:
        for layout in prs.slide_layouts:
            if pattern.lower() in layout.name.lower():
                return layout
    # Fallback to index
    if fallback_index < len(prs.slide_layouts):
        return prs.slide_layouts[fallback_index]
    return prs.slide_layouts[0]


def _find_layout_strict(prs, name_patterns):
    for pattern in name_patterns:
        for layout in prs.slide_layouts:
            if pattern.lower() in layout.name.lower():
                return layout
    names = ", ".join([l.name for l in prs.slide_layouts])
    raise SystemExit(f"Required slide layout not found. Searched patterns: {name_patterns}. Available layouts: {names}")


def get_cover_layout(prs):
    """Get the cover/title slide layout."""
    return _find_layout(prs, LAYOUT_COVER_NAMES, fallback_index=0)


def get_content_layout(prs):
    """Get the content slide layout (title + content area)."""
    return _find_layout(prs, LAYOUT_CONTENT_NAMES, fallback_index=1)


def get_title_only_layout(prs):
    """Get the title-only slide layout (enforced for all non-cover slides)."""
    return _find_layout_strict(prs, LAYOUT_TITLE_ONLY_NAMES)


def get_closing_layout(prs):
    """Get the closing/outro slide layout."""
    return _find_layout(prs, LAYOUT_CLOSING_NAMES, fallback_index=0)


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  HELPERS                                                          ║
# ╚═══════════════════════════════════════════════════════════════════╝

def _apply_theme_color(color_format, theme_color, brightness=0.0):
    """Apply a theme color with optional brightness adjustment."""
    color_format.theme_color = theme_color
    if brightness != 0.0:
        color_format.brightness = brightness


def set_bg(slide):
    """No-op: let the slide inherit its background from the slide master/layout.
    This ensures that when the user switches the PowerPoint theme, the
    background changes along with everything else. If you explicitly set
    a background here, it becomes an override that blocks theme inheritance.
    """
    pass


def txt(slide, l, t, w, h, text, sz=14, theme_color=MSO_THEME_COLOR.LIGHT_1,
        brightness=0.0, bold=False, align=PP_ALIGN.LEFT):
    """Add a positioned text box. Default color is LIGHT_1 (white on dark themes).
    For textboxes (not placeholders), PowerPoint does NOT inherit the
    master's text color, so we must set it explicitly."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    if theme_color is not None:
        _apply_theme_color(p.font.color, theme_color, brightness)
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return box


def rect(slide, l, t, w, h, fill_theme, fill_brightness=0.0,
         text="", sz=14, fc_theme=MSO_THEME_COLOR.LIGHT_1, fc_brightness=0.0,
         bold=False):
    """Add a rounded rectangle. Text defaults to LIGHT_1 (white)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    _apply_theme_color(shape.fill.fore_color, fill_theme, fill_brightness)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        if fc_theme is not None:
            _apply_theme_color(p.font.color, fc_theme, fc_brightness)
        p.font.bold = bold
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def header(slide, title, subtitle=""):
    """Set the slide title via the layout's title placeholder.
    Falls back to a floating text box if no title placeholder exists.
    Subtitle is always a floating text box positioned below the title.
    """
    set_bg(slide)
    # Try to use the layout's title placeholder (idx 0)
    title_set = False
    if slide.placeholders:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # Title placeholder
                ph.text = title
                title_set = True
                break
    if not title_set:
        txt(slide, 0.5, 0.3, 12, 0.6, title, sz=28, bold=True)
    if subtitle:
        txt(slide, 0.5, 0.85, 12, 0.4, subtitle, sz=14)


def card(slide, l, t, w, h, title, lines, accent_theme, tsz=16, bsz=14):
    """Add an info card with accent top border and bullet lines.
    All body lines are rendered as paragraphs inside a single text frame
    so they stay as one editable text block in PowerPoint.
    Card background is transparent — inherits from slide/theme background.
    A subtle border provides visual separation on any theme.
    """
    # Card background — transparent with accent border
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    bg.fill.background()  # transparent
    bg.line.width = Pt(1.5)
    _apply_theme_color(bg.line.color, accent_theme)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(l + 0.05), Inches(t + 0.05),
                                 Inches(w - 0.1), Inches(0.06))
    bar.fill.solid()
    _apply_theme_color(bar.fill.fore_color, accent_theme)
    bar.line.fill.background()
    # Title
    txt(slide, l + 0.15, t + 0.18, w - 0.3, 0.35, title,
        sz=tsz, theme_color=accent_theme, bold=True)
    # Body lines as a single text frame with multiple paragraphs
    body_box = slide.shapes.add_textbox(
        Inches(l + 0.15), Inches(t + 0.55),
        Inches(w - 0.3), Inches(h - 0.7)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(bsz)
        _apply_theme_color(p.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)


def phase_flow(slide, phases, y_top=1.5):
    """Draw a horizontal flow of numbered phase boxes with arrows.
    phases: list of (number, name, description, accent_theme_color)
    """
    x = 0.3
    for num, name, desc, clr in phases:
        rect(slide, x, y_top, 1.1, 0.5, clr, text=num, sz=16, bold=True)
        txt(slide, x, y_top + 0.55, 1.1, 0.35, name,
            sz=11, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, x, y_top + 0.9, 1.1, 0.6, desc,
            sz=9, align=PP_ALIGN.CENTER)
        if num != phases[-1][0]:
            txt(slide, x + 1.05, y_top + 0.05, 0.15, 0.4, "→",
                sz=16, align=PP_ALIGN.CENTER)
        x += 1.18


def grid_items(slide, items, y_start=1.5, cols=4):
    """Draw a grid of labeled items with colored accent dots.

    Each cell uses a dark-tinted accent fill so text remains readable on
    both light and dark themes.  Name and subtitle are grouped into a
    single text frame for easy editing in PowerPoint.

    items: list of (name, subtitle, accent_theme_color)
    """
    col_w = 3.0
    row_h = 0.9
    for i, (name, sub, clr) in enumerate(items):
        col = i % cols
        row = i // cols
        x = 0.5 + col * 3.2
        y = y_start + row * 1.1
        rect(slide, x, y, col_w, row_h, clr, 0.15)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(x + 0.12), Inches(y + 0.18),
                                      Inches(0.22), Inches(0.22))
        dot.fill.solid()
        _apply_theme_color(dot.fill.fore_color, clr)
        dot.line.fill.background()
        body = slide.shapes.add_textbox(
            Inches(x + 0.45), Inches(y + 0.1),
            Inches(col_w - 0.6), Inches(row_h - 0.2))
        tf = body.text_frame
        tf.word_wrap = True
        p_name = tf.paragraphs[0]
        p_name.text = name
        p_name.font.size = Pt(14)
        p_name.font.bold = True
        p_name.font.name = "Segoe UI"
        _apply_theme_color(p_name.font.color, MSO_THEME_COLOR.LIGHT_1)
        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.name = "Segoe UI"
        _apply_theme_color(p_sub.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
        p_sub.space_before = Pt(2)


def badge_row(slide, badges, y=5.0):
    """Draw a horizontal row of principle/feature badges.
    badges: list of (label, description, accent_theme_color)
    """
    n = len(badges)
    total_w = n * 2.1
    x = (13.333 - total_w) / 2
    for label, desc, clr in badges:
        rect(slide, x, y, 2.0, 0.45, clr, text=label, sz=10, bold=True)
        txt(slide, x, y + 0.5, 2.0, 0.35, desc, sz=10,
            align=PP_ALIGN.CENTER)
        x += 2.1


# ── Enhanced Visual Functions for Documentation ──

def create_enhanced_doc_bullets(slide, bullets, y_start=1.4):
    """Create enhanced bullet points for documentation with icons."""
    colors = [ACCENT_PRIMARY, ACCENT_INFO, ACCENT_POSITIVE, ACCENT_WARNING, ACCENT_AI]
    icons = ["●", "■", "▲", "◆", "★"]
    
    for i, bullet in enumerate(bullets[:5]):  # Max 5 for documentation
        y_pos = y_start + i * 0.6
        
        # Add icon circle
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(0.7), Inches(y_pos), 
                                        Inches(0.25), Inches(0.25))
        icon_bg.fill.solid()
        _apply_theme_color(icon_bg.fill.fore_color, colors[i % len(colors)])
        icon_bg.line.fill.background()
        
        # Add icon text
        txt(slide, 0.73, y_pos + 0.03, 0.16, 0.16, icons[i % len(icons)], 
            sz=10, theme_color=MSO_THEME_COLOR.LIGHT_1, bold=True, align=PP_ALIGN.CENTER)
        
        # Add bullet text with subtle background
        text_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(1.1), Inches(y_pos - 0.03), 
                                        Inches(11.2), Inches(0.35))
        text_bg.fill.solid()
        _apply_theme_color(text_bg.fill.fore_color, MSO_THEME_COLOR.LIGHT_2, -0.98)
        text_bg.line.fill.background()
        
        # Add bullet text
        txt(slide, 1.2, y_pos, 10.8, 0.3, bullet, sz=12)


def create_doc_feature_cards(slide, features, y_start=1.4):
    """Create professional feature cards for documentation."""
    card_width = 3.7
    card_height = 1.8
    colors = [ACCENT_PRIMARY, ACCENT_INFO, ACCENT_POSITIVE]
    
    for i, (title, description) in enumerate(features[:3]):
        x_pos = 0.7 + (i % 3) * 4.2
        y_pos = y_start + (i // 3) * 2.2
        
        # Create card with subtle gradient effect
        card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x_pos), Inches(y_pos), 
                                        Inches(card_width), Inches(card_height))
        card_bg.fill.solid()
        _apply_theme_color(card_bg.fill.fore_color, colors[i % len(colors)], -0.9)
        card_bg.line.width = Pt(1.5)
        _apply_theme_color(card_bg.line.color, colors[i % len(colors)])
        
        # Add icon header
        icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                            Inches(x_pos + 0.15), Inches(y_pos + 0.1), 
                                            Inches(0.35), Inches(0.35))
        icon_circle.fill.solid()
        _apply_theme_color(icon_circle.fill.fore_color, colors[i % len(colors)])
        icon_circle.line.fill.background()
        
        # Add feature title
        txt(slide, x_pos + 0.6, y_pos + 0.15, card_width - 0.9, 0.35, 
            title, sz=13, theme_color=colors[i % len(colors)], bold=True)
        
        # Add description
        txt(slide, x_pos + 0.15, y_pos + 0.55, card_width - 0.3, 1.0, 
            description, sz=11, theme_color=MSO_THEME_COLOR.LIGHT_1, brightness=-0.15)


def create_doc_process_flow(slide, steps, y_start=2.0):
    """Create a documentation-appropriate process flow."""
    box_width = 2.8
    box_height = 1.0
    arrow_length = 0.8
    
    for i, (step, desc) in enumerate(steps[:4]):
        x_pos = 1.2 + i * (box_width + arrow_length)
        y_pos = y_start
        
        # Process box
        process_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(x_pos), Inches(y_pos), 
                                            Inches(box_width), Inches(box_height))
        process_box.fill.solid()
        _apply_theme_color(process_box.fill.fore_color, ACCENT_INFO, -0.85)
        process_box.line.width = Pt(1.5)
        _apply_theme_color(process_box.line.color, ACCENT_INFO)
        
        # Step number
        number_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                              Inches(x_pos + 0.1), Inches(y_pos + 0.1), 
                                              Inches(0.35), Inches(0.35))
        number_circle.fill.solid()
        _apply_theme_color(number_circle.fill.fore_color, ACCENT_PRIMARY)
        number_circle.line.fill.background()
        
        txt(slide, x_pos + 0.2, y_pos + 0.18, 0.15, 0.15, str(i + 1), 
            sz=11, theme_color=MSO_THEME_COLOR.LIGHT_1, bold=True, align=PP_ALIGN.CENTER)
        
        # Step text
        txt(slide, x_pos + 0.55, y_pos + 0.2, box_width - 0.8, 0.25, 
            step, sz=11, theme_color=ACCENT_PRIMARY, bold=True)
        
        # Description
        txt(slide, x_pos + 0.2, y_pos + 0.5, box_width - 0.4, 0.35, 
            desc, sz=9)
        
        # Arrow to next step
        if i < len(steps) - 1 and i < 3:
            arrow_x = x_pos + box_width
            txt(slide, arrow_x, y_pos + 0.35, arrow_length, 0.25, "→", 
                sz=20, theme_color=ACCENT_POSITIVE, bold=True, align=PP_ALIGN.CENTER)


def create_doc_benefits_showcase(slide, benefits, y_start=1.4):
    """Create a documentation-focused benefits showcase."""
    # Add subtle header
    header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.5), Inches(y_start - 0.2), 
                                        Inches(12.3), Inches(0.4))
    header_rect.fill.solid()
    _apply_theme_color(header_rect.fill.fore_color, ACCENT_POSITIVE, -0.85)
    header_rect.line.fill.background()
    
    txt(slide, 0.7, y_start - 0.15, 12, 0.3, "Key Benefits", sz=14, 
        theme_color=ACCENT_POSITIVE, bold=True)
    
    # Create benefit items with subtle indicators
    for i, benefit in enumerate(benefits[:4]):
        y_pos = y_start + 0.4 + i * 0.7
        
        # Checkmark circle
        check_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                             Inches(0.7), Inches(y_pos), 
                                             Inches(0.25), Inches(0.25))
        check_circle.fill.solid()
        _apply_theme_color(check_circle.fill.fore_color, ACCENT_POSITIVE)
        check_circle.line.fill.background()
        
        # Checkmark
        txt(slide, 0.73, y_pos + 0.03, 0.16, 0.16, "✓", 
            sz=12, theme_color=MSO_THEME_COLOR.LIGHT_1, bold=True, align=PP_ALIGN.CENTER)
        
        # Benefit text
        txt(slide, 1.1, y_pos, 11, 0.3, benefit, sz=12)


def add_doc_visual_accents(slide):
    """Add subtle visual accents for documentation slides."""
    # Add corner accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                   Inches(0), Inches(0), 
                                   Inches(0.3), Inches(0.3))
    accent.fill.solid()
    _apply_theme_color(accent.fill.fore_color, ACCENT_PRIMARY, -0.8)
    accent.line.fill.background()
    
    # Add subtle footer line
    footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.5), Inches(6.7), 
                                        Inches(12.3), Inches(0.03))
    footer_line.fill.solid()
    _apply_theme_color(footer_line.fill.fore_color, ACCENT_INFO, -0.8)
    footer_line.line.fill.background()


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  DATA SECTION — replace with content from your docs               ║
# ╚═══════════════════════════════════════════════════════════════════╝

TITLE = "API Consistency Analysis"
INTENT = "Comprehensive API consistency analysis using code-mapper foundation for improved code quality"
SOURCE = "Data: distilled from analyze-api-consistency skill documentation (skill.md, description.md, tutorial.md)"
AI_NOTICE = "AI-generated presentation — content verified by author"

# Executive summary: 3 cards (problem / solution / output)
EXEC_PROBLEM = [
    "API inconsistencies confuse developers and break expectations",
    "Inconsistent naming patterns reduce code discoverability", 
    "Parameter variations create error-prone APIs",
    "Manual consistency analysis is time-consuming and error-prone"
]

EXEC_SOLUTION = [
    "Automated foundation analysis with code-mapper integration",
    "Systematic extraction of public API signatures",
    "Multi-dimensional consistency analysis (naming, parameters, returns)",
    "Structured reporting with severity classification"
]

EXEC_OUTPUT = [
    "Comprehensive API consistency report with actionable insights",
    "Severity-based inconsistency categorization (HIGH/MEDIUM/LOW)",
    "Canonical pattern recommendations and improvement roadmap",
    "Risk assessment and phased implementation plan"
]

# Process: 3 phases (validation / analysis / reporting)
PROCESS_PHASES = [
    ("Validation & Foundation", "Check inputs, validate project path, ensure full foundation outputs exist"),
    ("API Extraction & Analysis", "Extract public signatures, analyze naming/parameter/return patterns, identify inconsistencies"),
    "Report Generation & Recommendations"
]

# Component inventory: 4 key capabilities
COMPONENTS = [
    ("Foundation Integration", "Full code-mapper analysis with signature extraction", ACCENT_PRIMARY),
    ("Consistency Analysis", "Multi-dimensional pattern analysis with severity classification", ACCENT_INFO),
    ("Structured Reporting", "Template-based reports with actionable recommendations", ACCENT_POSITIVE),
    ("Improvement Roadmap", "Phased approach with risk assessment and rollback planning", ACCENT_WARNING)
]

# Usage example: 4 key steps
USAGE_STEPS = [
    "Run full foundation: python scripts/code-mapper/run.py --foundation <project-path>",
    "Invoke skill with project path and target modules",
    "Review extracted API signatures and consistency analysis", 
    "Save structured report with improvement recommendations"
]

# Success criteria: 4 key metrics
SUCCESS_METRICS = [
    ("Complete Analysis", "All requested modules analyzed with full API coverage"),
    ("Inconsistency Detection", "Major inconsistencies identified and categorized by severity"),
    ("Canonical Patterns", "Recommended standard patterns for repository consistency"),
    ("Actionable Roadmap", "Phased improvement plan with risk assessment")
]

SOURCES = [
    ("HAAL Skills Repository", "https://github.com/haal-ai/haal-skills"),
    ("Skill Documentation", "skills/analyze-api-consistency/"),
    ("Code-Mapper Foundation", "scripts/code-mapper/"),
]

TEMPLATE_PATH = SCRIPT_DIR / "template.pptx"

# Output path
OUTPUT_PATH = SCRIPT_DIR / f"{SCRIPT_DIR.name}.pptx"

# ... (rest of the code remains the same)
# ║  SLIDE BUILDERS — adapt these to your content                     ║
# ╚═══════════════════════════════════════════════════════════════════╝

def build_title_slide(prs):
    slide = prs.slides.add_slide(get_cover_layout(prs))
    set_bg(slide)
    # Populate the cover layout's title placeholder
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title
            ph.text = TITLE
        elif idx == 13:  # Subtitle (template-specific)
            ph.text = INTENT
        elif idx == 18:  # Lower body area (template-specific)
            ph.text = f"{SOURCE}\n{AI_NOTICE}"


def build_exec_summary(prs):
    slide = prs.slides.add_slide(get_title_only_layout(prs))
    header(slide, "Executive Summary")
    add_doc_visual_accents(slide)
    
    # Use enhanced feature cards instead of regular cards
    features = [
        ("Problem", "API inconsistencies confuse developers and reduce code quality"),
        ("Solution", "Automated foundation analysis with systematic consistency checking"),
        ("Output", "Structured reports with actionable improvement recommendations")
    ]
    create_doc_feature_cards(slide, features, y_start=1.4)
    
    # Add insight box
    rect(slide, 1.5, 5.0, 10.3, 0.55,
         ACCENT_PRIMARY, fill_brightness=-0.7,
         text="Systematic approach eliminates manual consistency review errors", sz=14, fc_theme=ACCENT_POSITIVE)


def build_process_overview(prs):
    """Enhanced process overview with visual flow."""
    slide = prs.slides.add_slide(get_title_only_layout(prs))
    header(slide, "Process Overview", "Three-phase workflow for API consistency analysis")
    add_doc_visual_accents(slide)
    
    # Use enhanced process flow
    steps = [
        ("Validation & Foundation", "Check inputs, validate project path, ensure foundation outputs"),
        ("API Extraction & Analysis", "Extract signatures, analyze patterns, identify inconsistencies"),
        ("Report Generation", "Create structured report with recommendations and roadmap")
    ]
    create_doc_process_flow(slide, steps, y_start=2.0)
    
    # Add key benefits
    benefits = [
        "Comprehensive coverage of all public APIs",
        "Multi-dimensional consistency analysis (naming, parameters, returns)",
        "Severity-based prioritization for focused improvements"
    ]
    create_doc_benefits_showcase(slide, benefits, y_start=3.5)


def build_component_inventory(prs):
    """Enhanced component inventory with visual cards."""
    slide = prs.slides.add_slide(get_title_only_layout(prs))
    header(slide, "Core Capabilities", "Key features of the API consistency analysis system")
    add_doc_visual_accents(slide)
    
    # Use enhanced feature cards for components
    features = [
        ("Foundation Integration", "Full code-mapper analysis with complete signature extraction"),
        ("Consistency Analysis", "Multi-dimensional pattern analysis with severity classification"),
        ("Structured Reporting", "Template-based reports with actionable recommendations"),
        ("Improvement Roadmap", "Phased approach with risk assessment and rollback planning")
    ]
    create_doc_feature_cards(slide, features, y_start=1.4)


def build_usage_example(prs):
    """Enhanced usage example with visual steps."""
    slide = prs.slides.add_slide(get_title_only_layout(prs))
    header(slide, "Usage Example", "Step-by-step workflow for API consistency analysis")
    add_doc_visual_accents(slide)
    
    # Use enhanced bullets for usage steps
    usage_steps = [
        "Run full foundation: python scripts/code-mapper/run.py --foundation <project-path>",
        "Invoke skill with project path and target modules to analyze",
        "Review extracted API signatures and detailed consistency analysis", 
        "Save structured report with improvement recommendations and roadmap"
    ]
    create_enhanced_doc_bullets(slide, usage_steps, y_start=1.8)
    
    # Add success metrics
    card(slide, 0.5, 4.5, 12.3, 1.8, "Success Metrics", [
        "Complete analysis of all requested modules with full API coverage",
        "Major inconsistencies identified and categorized by severity level",
        "Canonical patterns recommended for repository consistency",
        "Actionable roadmap with risk assessment and implementation phases"
    ], ACCENT_POSITIVE, bsz=12)


def build_sources_slide(prs):
    """Final slide: Sources & References with clickable hyperlinks."""
    if not SOURCES:
        return
    slide = prs.slides.add_slide(get_title_only_layout(prs))
    header(slide, "Sources & References")

    body = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5)
    )
    tf = body.text_frame
    tf.word_wrap = True
    for i, (label, url) in enumerate(SOURCES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        # Bullet prefix as plain run
        run_bullet = p.add_run()
        run_bullet.text = "• "
        run_bullet.font.size = Pt(14)
        _apply_theme_color(run_bullet.font.color, MSO_THEME_COLOR.LIGHT_1, -0.25)
        run_bullet.font.name = "Segoe UI"
        # Label as clickable hyperlink run
        run_link = p.add_run()
        run_link.text = label
        run_link.font.size = Pt(14)
        run_link.font.name = "Segoe UI"
        _apply_theme_color(run_link.font.color, MSO_THEME_COLOR.HYPERLINK)
        run_link.font.underline = True
        run_link.hyperlink.address = url
        # URL shown in muted text after the label
        run_url = p.add_run()
        run_url.text = f"  ({url})"
        run_url.font.size = Pt(10)
        run_url.font.name = "Segoe UI"
        _apply_theme_color(run_url.font.color, MSO_THEME_COLOR.LIGHT_2)


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  MAIN                                                             ║
# ╚═══════════════════════════════════════════════════════════════════╝

def main():
    _ensure_safe_execution_location()
    # Use template if provided, otherwise blank presentation
    if TEMPLATE_PATH and Path(TEMPLATE_PATH).exists():
        prs = Presentation(str(TEMPLATE_PATH))
        print(f"📎 Using template: {TEMPLATE_PATH}")
    else:
        prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)
    build_exec_summary(prs)
    build_process_overview(prs)
    build_component_inventory(prs)
    build_usage_example(prs)
    build_sources_slide(prs)

    # Set document properties
    prs.core_properties.title = TITLE
    prs.core_properties.subject = INTENT
    prs.core_properties.comments = "Generated by distill-docs-to-pptx skill"

    prs.save(str(OUTPUT_PATH))
    print(f"✅ Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
